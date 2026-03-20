import io
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


def extract_docx_text(file_path: str | Path) -> str:
    from docx import Document

    doc = Document(str(file_path))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            if para.style and para.style.name.startswith("Heading"):
                level = para.style.name.replace("Heading ", "").strip()
                try:
                    level = int(level)
                except ValueError:
                    level = 1
                paragraphs.append(f"{'#' * level} {para.text}")
            else:
                paragraphs.append(para.text)
    return "\n\n".join(paragraphs)


def extract_pptx_text(file_path: str | Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(file_path))
    slides_text = []
    for i, slide in enumerate(prs.slides, 1):
        slide_parts = [f"## Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
        slides_text.append("\n\n".join(slide_parts))
    return "\n\n---\n\n".join(slides_text)


def extract_pdf_pages_as_images(
    file_path: str | Path,
    dpi: int = 150,
    max_dimension: int = 1536,
    jpeg_quality: int = 85,
) -> list[bytes]:
    from pdf2image import convert_from_path

    images = convert_from_path(str(file_path), dpi=dpi)
    result = []
    for img in images:
        # Downscale if either dimension exceeds max_dimension
        w, h = img.size
        if w > max_dimension or h > max_dimension:
            scale = max_dimension / max(w, h)
            img = img.resize((int(w * scale), int(h * scale)), resample=1)  # LANCZOS
        # Convert to RGB (JPEG doesn't support alpha)
        if img.mode in ("RGBA", "P"):
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=jpeg_quality)
        result.append(buf.getvalue())
    return result


def extract_pdf_text_fast(file_path: str | Path, dpi: int = 200) -> str:
    """Extract text from PDF using Tesseract OCR. Fast but less context-aware."""
    import pytesseract
    from pdf2image import convert_from_path

    images = convert_from_path(str(file_path), dpi=dpi)
    pages = []
    for i, img in enumerate(images, 1):
        text = pytesseract.image_to_string(img).strip()
        pages.append(f"## Page {i}\n\n{text}")
    return "\n\n---\n\n".join(pages)


def extract_image_text_fast(file_path: str | Path) -> str:
    """Extract text from a single image using Tesseract OCR."""
    import pytesseract
    from PIL import Image

    img = Image.open(str(file_path))
    text = pytesseract.image_to_string(img).strip()
    return text


def detect_file_type(filename: str) -> str:
    ext = Path(filename).suffix.lower()
    type_map = {
        ".pdf": "pdf",
        ".docx": "docx",
        ".doc": "docx",
        ".pptx": "pptx",
        ".ppt": "pptx",
        ".png": "image",
        ".jpg": "image",
        ".jpeg": "image",
        ".gif": "image",
        ".bmp": "image",
        ".tiff": "image",
        ".webp": "image",
    }
    return type_map.get(ext, "pdf")
